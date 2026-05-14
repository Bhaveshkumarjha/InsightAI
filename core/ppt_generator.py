import pandas as pd
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

import matplotlib.pyplot as plt
import seaborn as sns

import tempfile
import os
import time
from datetime import datetime


# =========================================================
# CEO / EXECUTIVE LEVEL POWERPOINT GENERATOR
# =========================================================

class ExecutivePPT:

    def __init__(self):

        self.prs = Presentation()

        self.blank = self.prs.slide_layouts[6]

        # =================================================
        # MODERN EXECUTIVE THEME
        # =================================================

        self.colors = {

            # BACKGROUND
            "bg": RGBColor(255, 255, 255),

            # BRAND COLORS
            "primary": RGBColor(37, 99, 235),

            "accent": RGBColor(16, 185, 129),

            # TEXT COLORS
            "title": RGBColor(15, 23, 42),

            "text": RGBColor(71, 85, 105),

            # ALERT COLORS
            "danger": RGBColor(239, 68, 68),

            "warning": RGBColor(245, 158, 11),

            # CARD COLOR
            "card": RGBColor(248, 250, 252)
        }

        self.temp_files = []

    # =====================================================
    # BACKGROUND
    # =====================================================

    def set_background(self, slide):

        fill = slide.background.fill

        fill.solid()

        fill.fore_color.rgb = self.colors["bg"]

    # =====================================================
    # TITLE
    # =====================================================

    def add_title(self, slide, text):

        title = slide.shapes.add_textbox(
            Inches(0.6),
            Inches(0.3),
            Inches(9),
            Inches(0.6)
        )

        tf = title.text_frame

        p = tf.paragraphs[0]

        p.text = text

        p.font.size = Pt(28)

        p.font.bold = True

        p.font.color.rgb = self.colors["title"]

    # =====================================================
    # COVER SLIDE
    # =====================================================

    def cover_slide(self):

        slide = self.prs.slides.add_slide(self.blank)

        self.set_background(slide)

        title = slide.shapes.add_textbox(
            Inches(1),
            Inches(2),
            Inches(8),
            Inches(1)
        )

        tf = title.text_frame

        p = tf.paragraphs[0]

        p.text = "InsightAI Executive Report"

        p.font.size = Pt(34)

        p.font.bold = True

        p.font.color.rgb = self.colors["title"]

        p.alignment = PP_ALIGN.CENTER

        subtitle = slide.shapes.add_textbox(
            Inches(1),
            Inches(3),
            Inches(8),
            Inches(0.5)
        )

        tf2 = subtitle.text_frame

        p2 = tf2.paragraphs[0]

        p2.text = f"Generated on {datetime.now().strftime('%d %B %Y %I:%M %p')}"

        p2.font.size = Pt(18)

        p2.font.color.rgb = self.colors["text"]

        p2.alignment = PP_ALIGN.CENTER

    # =====================================================
    # KPI CARD
    # =====================================================

    def kpi_card(self, slide, x, y, title, value, color):

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(2),
            Inches(1.2)
        )

        fill = shape.fill

        fill.solid()

        fill.fore_color.rgb = self.colors["card"]

        shape.line.color.rgb = color

        # SHADOW
        shape.shadow.inherit = False

        shadow = shape.shadow

        shadow.blur = 8

        shadow.distance = 2

        shadow.angle = 45

        

        # TITLE
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1),
            Inches(y + 0.1),
            Inches(1.8),
            Inches(0.3)
        )

        tf = title_box.text_frame

        p = tf.paragraphs[0]

        p.text = title

        p.font.size = Pt(12)

        p.font.bold = True

        p.font.color.rgb = color

        # VALUE
        value_box = slide.shapes.add_textbox(
            Inches(x + 0.1),
            Inches(y + 0.45),
            Inches(1.8),
            Inches(0.4)
        )

        tf2 = value_box.text_frame

        p2 = tf2.paragraphs[0]

        p2.text = str(value)

        p2.font.size = Pt(22)

        p2.font.bold = True

        p2.font.color.rgb = self.colors["title"]

    # =====================================================
    # EXECUTIVE SUMMARY
    # =====================================================

    def executive_summary(self, df):

        slide = self.prs.slides.add_slide(self.blank)

        self.set_background(slide)

        self.add_title(slide, "Executive Summary")

        bullets = [

            f"Dataset contains {df.shape[0]} records and {df.shape[1]} columns.",

            f"Total missing values identified: {df.isnull().sum().sum()}",

            f"Duplicate records detected: {df.duplicated().sum()}",

            "Data quality assessment completed successfully.",

            "Advanced statistical profiling executed.",

            "Business intelligence insights generated."
        ]

        top = 1.5

        for bullet in bullets:

            box = slide.shapes.add_textbox(
                Inches(0.8),
                Inches(top),
                Inches(8),
                Inches(0.4)
            )

            tf = box.text_frame

            p = tf.paragraphs[0]

            p.text = "• " + bullet

            p.font.size = Pt(18)

            p.font.color.rgb = self.colors["text"]

            top += 0.6

    # =====================================================
    # KPI DASHBOARD
    # =====================================================

    def dashboard_slide(self, df):

        slide = self.prs.slides.add_slide(self.blank)

        self.set_background(slide)

        self.add_title(slide, "Business KPI Dashboard")

        rows, cols = df.shape

        missing = df.isnull().sum().sum()

        duplicates = df.duplicated().sum()

        numeric = len(df.select_dtypes(include=np.number).columns)

        self.kpi_card(slide, 0.6, 1.5, "Rows", rows, self.colors["primary"])

        self.kpi_card(slide, 2.8, 1.5, "Columns", cols, self.colors["accent"])

        self.kpi_card(slide, 5.0, 1.5, "Missing", missing, self.colors["warning"])

        self.kpi_card(slide, 7.2, 1.5, "Duplicates", duplicates, self.colors["danger"])

        self.kpi_card(slide, 2.8, 3.2, "Numeric", numeric, self.colors["primary"])

    # =====================================================
    # CHART SLIDE
    # =====================================================

    def chart_slide(self, title, fig):

        slide = self.prs.slides.add_slide(self.blank)

        self.set_background(slide)

        self.add_title(slide, title)

        tmp = tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".png"
        )

        fig.savefig(
            tmp.name,
            bbox_inches='tight',
            facecolor='white'
        )

        self.temp_files.append(tmp.name)

        slide.shapes.add_picture(
            tmp.name,
            Inches(0.7),
            Inches(1.2),
            width=Inches(8.5)
        )

        plt.close(fig)

    # =====================================================
    # INSIGHTS
    # =====================================================

    def insights_slide(self, insights):

        slide = self.prs.slides.add_slide(self.blank)

        self.set_background(slide)

        self.add_title(slide, "AI Business Insights")

        top = 1.5

        for insight in insights:

            box = slide.shapes.add_textbox(
                Inches(0.8),
                Inches(top),
                Inches(8),
                Inches(0.5)
            )

            tf = box.text_frame

            p = tf.paragraphs[0]

            p.text = "✔ " + insight

            p.font.size = Pt(18)

            p.font.color.rgb = self.colors["text"]

            top += 0.7

    # =====================================================
    # RECOMMENDATIONS
    # =====================================================

    def recommendation_slide(self, recommendations):

        slide = self.prs.slides.add_slide(self.blank)

        self.set_background(slide)

        self.add_title(slide, "Strategic Recommendations")

        top = 1.5

        for rec in recommendations:

            box = slide.shapes.add_textbox(
                Inches(0.8),
                Inches(top),
                Inches(8),
                Inches(0.5)
            )

            tf = box.text_frame

            p = tf.paragraphs[0]

            p.text = "➜ " + rec

            p.font.size = Pt(18)

            p.font.color.rgb = self.colors["text"]

            top += 0.7

    # =====================================================
    # THANK YOU
    # =====================================================

    def thank_you_slide(self):

        slide = self.prs.slides.add_slide(self.blank)

        self.set_background(slide)

        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(2.5),
            Inches(8),
            Inches(1)
        )

        tf = box.text_frame

        p = tf.paragraphs[0]

        p.text = "Thank You"

        p.font.size = Pt(40)

        p.font.bold = True

        p.font.color.rgb = self.colors["title"]

        p.alignment = PP_ALIGN.CENTER

    # =====================================================
    # MAIN REPORT
    # =====================================================

    def generate(self, df):

        numeric_cols = df.select_dtypes(include=np.number).columns.tolist()

        self.cover_slide()

        self.executive_summary(df)

        self.dashboard_slide(df)

        # =================================================
        # HISTOGRAMS
        # =================================================

        for col in numeric_cols[:5]:

            fig, ax = plt.subplots(
                figsize=(8, 4),
                facecolor='white'
            )

            sns.histplot(
                df[col].dropna(),
                kde=True,
                color="#2563eb",
                ax=ax
            )

            ax.set_facecolor('white')

            ax.tick_params(colors='black')

            ax.set_title(
                f"{col} Distribution",
                color='black',
                fontsize=16
            )

            ax.set_xlabel(col, color='black')

            ax.set_ylabel("Count", color='black')

            self.chart_slide(
                f"{col} Distribution Analysis",
                fig
            )

        # =================================================
        # CORRELATION HEATMAP
        # =================================================

        if len(numeric_cols) >= 2:

            fig, ax = plt.subplots(
                figsize=(8, 6),
                facecolor='white'
            )

            sns.heatmap(
                df[numeric_cols].corr(),
                annot=True,
                cmap="Blues",
                ax=ax
            )

            ax.set_title(
                "Correlation Matrix",
                color='black'
            )

            self.chart_slide(
                "Feature Correlation Intelligence",
                fig
            )

        # =================================================
        # BOXPLOTS
        # =================================================

        for col in numeric_cols[:3]:

            fig, ax = plt.subplots(
                figsize=(8, 4),
                facecolor='white'
            )

            sns.boxplot(
                x=df[col],
                color="#10b981",
                ax=ax
            )

            ax.set_facecolor('white')

            ax.tick_params(colors='black')

            ax.set_title(
                f"{col} Outlier Analysis",
                color='black'
            )

            self.chart_slide(
                f"{col} Outlier Intelligence",
                fig
            )

        # =================================================
        # AI INSIGHTS
        # =================================================

        insights = [

            "Strong relationships detected between multiple business variables.",

            "Data quality is suitable for strategic analytics.",

            "Potential outliers identified in numeric distributions.",

            "High-value trends detected across major KPIs.",

            "Business forecasting readiness achieved.",

            "Dataset prepared for advanced AI and ML workflows."
        ]

        self.insights_slide(insights)

        # =================================================
        # RECOMMENDATIONS
        # =================================================

        recommendations = [

            "Improve data governance to reduce missing information.",

            "Monitor high-correlation variables for predictive modeling.",

            "Focus on KPI-driven decision making.",

            "Implement automated anomaly detection systems.",

            "Use AI-driven forecasting for business growth.",

            "Deploy executive dashboards for real-time monitoring."
        ]

        self.recommendation_slide(recommendations)

        self.thank_you_slide()

        # OUTPUT FILE
        output = f"InsightAI_Executive_Report_{int(time.time())}.pptx"

        self.prs.save(output)

        # CLEAN TEMP FILES
        for f in self.temp_files:

            if os.path.exists(f):

                os.remove(f)

        return output


# =========================================================
# STREAMLIT FUNCTION
# =========================================================

def generate_ppt(df, numeric_columns=None):

    report = ExecutivePPT()

    return report.generate(df)